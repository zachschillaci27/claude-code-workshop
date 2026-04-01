"""Generate PPTX slides from the workshop slide content."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# Theme colors
BG_DARK = RGBColor(0x0F, 0x0F, 0x1E)
BG_SLIDE = RGBColor(0x1A, 0x1A, 0x2E)
TEXT_MAIN = RGBColor(0xE0, 0xE0, 0xE0)
TEXT_HEADING = RGBColor(0xD4, 0xA5, 0x74)
TEXT_SUB = RGBColor(0xC4, 0x9A, 0x6C)
TEXT_DIM = RGBColor(0xB0, 0xB0, 0xB0)
TEXT_CODE = RGBColor(0xE8, 0xC4, 0x7C)
TEXT_ACCENT = RGBColor(0x7E, 0xB8, 0xDA)
BG_CODE = RGBColor(0x16, 0x16, 0x2A)
BG_TABLE_HDR = RGBColor(0x2D, 0x2D, 0x44)
BG_TABLE_ROW = RGBColor(0x1E, 0x1E, 0x36)

FONT_BODY = "Calibri"
FONT_CODE = "Consolas"


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    return txBox.text_frame


def add_title(tf, text, size=Pt(36), color=TEXT_HEADING):
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = size
    p.font.color.rgb = color
    p.font.bold = True
    p.font.name = FONT_BODY
    return p


def add_para(tf, text, size=Pt(18), color=TEXT_MAIN, bold=False, space_before=Pt(6), bullet=False, font_name=FONT_BODY):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = size
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.space_before = space_before
    if bullet:
        p.level = 0
    return p


def add_code_block(tf, lines, size=Pt(14)):
    for line in lines:
        p = tf.add_paragraph()
        p.text = line
        p.font.size = size
        p.font.color.rgb = TEXT_CODE
        p.font.name = FONT_CODE
        p.space_before = Pt(2)


def make_title_slide(prs, title, subtitle="", is_section=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BG_DARK)

    tf = add_text_box(slide, Inches(1), Inches(2), Inches(8), Inches(2))
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    add_title(tf, title, size=Pt(44) if not is_section else Pt(40))
    if subtitle:
        add_para(tf, subtitle, size=Pt(24), color=TEXT_SUB, space_before=Pt(16))
        tf.paragraphs[-1].alignment = PP_ALIGN.CENTER


def make_content_slide(prs, title, bullets, code=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_SLIDE)

    tf = add_text_box(slide, Inches(0.6), Inches(0.4), Inches(8.8), Inches(0.8))
    add_title(tf, title, size=Pt(32))

    y = Inches(1.3)
    if bullets:
        tf2 = add_text_box(slide, Inches(0.6), y, Inches(8.8), Inches(3.5) if code else Inches(5.5))
        first = True
        for b in bullets:
            if first:
                tf2.paragraphs[0].text = b
                tf2.paragraphs[0].font.size = Pt(18)
                tf2.paragraphs[0].font.color.rgb = TEXT_MAIN
                tf2.paragraphs[0].font.name = FONT_BODY
                tf2.paragraphs[0].space_before = Pt(4)
                first = False
            else:
                bold = b.startswith("**")
                text = b.strip("*") if bold else b
                add_para(tf2, text, bold=bold, space_before=Pt(8))
        y = Inches(4.8) if code else y

    if code:
        tf3 = add_text_box(slide, Inches(0.6), y, Inches(8.8), Inches(2.5))
        add_code_block(tf3, code)


def make_table_slide(prs, title, headers, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG_SLIDE)

    tf = add_text_box(slide, Inches(0.6), Inches(0.4), Inches(8.8), Inches(0.8))
    add_title(tf, title, size=Pt(32))

    cols = len(headers)
    n_rows = len(rows) + 1
    col_width = Inches(8.8 / cols)
    table_shape = slide.shapes.add_table(n_rows, cols, Inches(0.6), Inches(1.4), Inches(8.8), Inches(0.45 * n_rows))
    table = table_shape.table

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = BG_TABLE_HDR
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14)
            p.font.color.rgb = TEXT_HEADING
            p.font.bold = True
            p.font.name = FONT_BODY

    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = BG_TABLE_ROW
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(13)
                p.font.color.rgb = TEXT_MAIN
                p.font.name = FONT_BODY


def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- SLIDE 1: Title ---
    make_title_slide(prs, "Claude Code Workshop", "From Beginner to Intermediate\n\nAn interactive, hands-on session with a real codebase\n\n~85 minutes  |  8 exercises  |  Live coding")

    # --- SLIDE 2: Agenda ---
    make_table_slide(prs, "Agenda",
        ["Part", "Topics", "Time"],
        [
            ["1. Foundations", "Launch, explore, edit, recover", "25 min"],
            ["2. Configuration", "CLAUDE.md, settings, permissions", "20 min"],
            ["3. Automation", "Hooks, custom skills", "20 min"],
            ["4. Advanced", "Subagents, context, real-world workflow", "15 min"],
            ["Wrap-up", "Quick wins, Q&A", "5 min"],
        ])

    # --- SLIDE 3: Setup ---
    make_content_slide(prs, "Setup Check", [], code=[
        "claude --version         # Claude Code installed",
        "python --version         # Python 3.11+",
        "uv --version             # uv package manager",
        "",
        "git clone https://github.com/zachschillaci27/claude-code-workshop.git",
        "cd claude-code-workshop",
        "uv sync",
        "uv run pytest            # 18 tests should pass",
    ])

    # --- SLIDE 4: Part 1 section ---
    make_title_slide(prs, "Part 1", "Foundations", is_section=True)

    # --- SLIDE 5: What is Claude Code ---
    make_content_slide(prs, "Exercise 1 — Getting Started",
        [
            "What is Claude Code?",
            "",
            "An agentic coding assistant in your terminal",
            "Not autocomplete — it takes multi-step actions",
            "Reads, searches, edits files, runs commands",
            "You approve each action (unless configured otherwise)",
        ],
        code=["claude"])

    # --- SLIDE 6: Core Tools ---
    make_table_slide(prs, "Core Tools",
        ["Tool", "Purpose"],
        [
            ["Read", "Read file contents"],
            ["Edit", "Modify existing files"],
            ["Write", "Create new files"],
            ["Bash", "Run shell commands"],
            ["Grep", "Search file contents"],
            ["Glob", "Find files by pattern"],
        ])

    # --- SLIDE 7: Try It ---
    make_content_slide(prs, "Try It",
        [
            "Watch how Claude reads first, then edits, and asks for permission.",
        ],
        code=[
            '"What does this project do?"',
            "",
            '"Show me all the API endpoints"',
            "",
            '"Run the tests"',
            "",
            '"Add a description field to the /health endpoint"',
        ])

    # --- SLIDE 8: Essential Navigation ---
    make_content_slide(prs, "Essential Navigation",
        [
            "@-mentions — Point Claude at a file directly:",
            '    "Look at @src/taskflow/models.py and add a due_date field"',
            "",
            "Interrupt — Ctrl+C stops Claude mid-action",
            "",
            "Multi-line — Shift+Enter for longer prompts",
            "",
            "Permission modes — Shift+Tab cycles through:",
            "    default -> acceptEdits -> plan -> ...",
            "",
            "Plan mode = read-only. Claude designs but doesn't change anything.",
        ])

    # --- SLIDE 9: Context Management ---
    make_table_slide(prs, "Context Management",
        ["Command", "What it does"],
        [
            ["/compact", "Summarize conversation, free up space"],
            ["claude --continue", "Resume your last session"],
            ["claude --resume", "Pick from previous sessions"],
        ])

    # --- SLIDE 10: Part 2 section ---
    make_title_slide(prs, "Part 2", "Configuration", is_section=True)

    # --- SLIDE 11: CLAUDE.md ---
    make_content_slide(prs, "Exercise 2 — CLAUDE.md",
        [
            "The single most impactful file for productivity",
            "",
            "Claude reads it at the start of every session. It tells Claude:",
            "  - How to build and test your project",
            "  - Code conventions to follow",
            "  - Architecture overview",
            "  - Common workflows",
            "",
            "File locations:",
            "  ./CLAUDE.md              — Project (shared via git)",
            "  ./.claude/CLAUDE.md      — Alt project location",
            "  ~/.claude/CLAUDE.md      — User (personal, all projects)",
        ])

    # --- SLIDE 12: CLAUDE.md content ---
    make_content_slide(prs, "What to Put in CLAUDE.md", [], code=[
        "# TaskFlow API",
        "",
        "## Build & Run",
        "- Install: `uv sync`",
        "- Run tests: `uv run pytest`",
        "- Lint: `uv run ruff check src/ tests/`",
        "",
        "## Code Conventions",
        "- Use type hints on all function signatures",
        "- Prefer `str | None` over `Optional[str]`",
        "- Keep endpoint handlers thin",
        "",
        "## Architecture",
        "- src/taskflow/main.py  — FastAPI entry point",
        "- src/taskflow/models.py — Pydantic models",
    ])

    # --- SLIDE 13: CLAUDE.md effect ---
    make_content_slide(prs, "The Effect",
        [
            "Without CLAUDE.md:",
            "  Claude guesses at conventions, might use wrong commands",
            "",
            "With CLAUDE.md:",
            "  Claude follows your team's standards automatically",
            "",
            "Try it:",
            '  "Create a new endpoint to get tasks by tag"',
            "  Check: Type hints? /api/v1/ prefix? 404 handling?",
            "",
            "Bootstrap with /init — Claude analyzes your codebase and generates one",
        ])

    # --- SLIDE 14: Settings ---
    make_content_slide(prs, "Exercise 3 — Settings & Permissions",
        [
            "Settings enforce behavior. CLAUDE.md guides behavior.",
            "",
            ".claude/settings.json        — Project (shared via git)",
            ".claude/settings.local.json  — Local (gitignored)",
            "~/.claude/settings.json      — User (personal)",
            "",
            "The mental model:",
            "  allow = auto-approve (no permission prompt)",
            "  deny  = block entirely (Claude gets error feedback)",
            "  neither = prompt the user each time",
        ])

    # --- SLIDE 15: Permission rules ---
    make_content_slide(prs, "Permission Rules", [
            'Pattern syntax: ToolName(glob pattern)',
        ],
        code=[
            '{',
            '  "permissions": {',
            '    "allow": [',
            '      "Read", "Glob", "Grep",',
            '      "Bash(uv *)", "Bash(git *)"',
            '    ],',
            '    "deny": [',
            '      "Bash(rm -rf *)",',
            '      "Edit(.env*)"',
            '    ]',
            '  }',
            '}',
        ])

    # --- SLIDE 16: Part 3 section ---
    make_title_slide(prs, "Part 3", "Automation", is_section=True)

    # --- SLIDE 17: Hooks ---
    make_table_slide(prs, "Exercise 4 — Hooks",
        ["Event", "When", "Use For"],
        [
            ["PreToolUse", "Before a tool runs", "Block dangerous actions"],
            ["PostToolUse", "After a tool succeeds", "Auto-format, lint"],
            ["Notification", "Claude needs attention", "Desktop alerts"],
            ["SessionStart", "Session begins", "Load env vars"],
        ])

    # --- SLIDE 18: Hook examples ---
    make_content_slide(prs, "Hook Examples in This Project",
        [
            "Auto-format after every edit (PostToolUse):",
            "  ruff format + ruff check --fix runs automatically",
            "",
            "Block hardcoded secrets (PreToolUse):",
            "  .claude/hooks/check-secrets.sh checks for password/token patterns",
            "  Exit code 2 = block, stderr becomes Claude's feedback",
            "",
            "Hook I/O:",
            "  Input: JSON on stdin with tool_name, tool_input, session_id",
            "  Exit 0 = proceed  |  Exit 2 = block",
        ])

    # --- SLIDE 19: Demo hooks ---
    make_content_slide(prs, "Demo: Hooks in Action",
        [
            "Auto-format:",
        ],
        code=[
            '"Add a function to utils.py with really bad formatting"',
            "# -> ruff fixes it automatically after Claude writes",
            "",
            '"Add DATABASE_URL = \'postgresql://admin:secret@prod:5432/db\'"',
            "# -> Hook blocks the write! Claude gets feedback.",
        ])

    # --- SLIDE 20: Skills ---
    make_content_slide(prs, "Exercise 5 — Custom Skills",
        [
            "Reusable workflows as /slash-commands",
            "",
            ".claude/skills/<name>/SKILL.md    — Project",
            "~/.claude/skills/<name>/SKILL.md  — User",
        ],
        code=[
            "---",
            "name: review",
            'description: Review code for quality and security',
            'argument-hint: "[file-or-directory]"',
            "allowed-tools: Read, Grep, Glob",
            "model: sonnet",
            "effort: high",
            "---",
            "",
            "Review $ARGUMENTS against the project's coding standards.",
        ])

    # --- SLIDE 21: Skill fields ---
    make_table_slide(prs, "Skill Frontmatter Fields",
        ["Field", "Purpose"],
        [
            ["name", "Slash command identifier"],
            ["description", "When to use (shown in / menu)"],
            ["argument-hint", 'Autocomplete hint, e.g. "[filename]"'],
            ["allowed-tools", "Restrict tool access"],
            ["model / effort", "Cost/quality tradeoffs"],
            ["disable-model-invocation", "true = pure shell, no AI"],
            ["user-invocable", "false = hidden from menu"],
            ["context: fork + agent", "Delegate to a subagent"],
            ["paths", "Auto-load only for matching files"],
        ])

    # --- SLIDE 22: Try skills ---
    make_content_slide(prs, "Try the Pre-Built Skills", [],
        code=[
            "/review src/taskflow/routers/tasks.py",
            "# -> Structured code review against project standards",
            "",
            "/test-coverage",
            "# -> Finds untested functions, suggests test cases",
            "",
            '/add-endpoint "GET /api/v1/tasks/search - search by title"',
            "# -> Scaffolds endpoint + database method + tests",
        ])

    # --- SLIDE 23: Part 4 section ---
    make_title_slide(prs, "Part 4", "Advanced", is_section=True)

    # --- SLIDE 24: Subagents ---
    make_table_slide(prs, "Exercise 6 — Subagents",
        ["Built-in", "Model", "Tools", "Use For"],
        [
            ["Explore", "Haiku (fast)", "Read-only", "Search, analyze"],
            ["Plan", "Inherited", "Read-only", "Design strategy"],
            ["general-purpose", "Inherited", "All", "Complex tasks"],
        ])

    # --- SLIDE 25: Agent fields ---
    make_table_slide(prs, "Agent Frontmatter Fields",
        ["Field", "Purpose"],
        [
            ["name / description", "Identity (required)"],
            ["tools / disallowedTools", "Allow/deny tools"],
            ["model", "haiku, sonnet, opus, inherit"],
            ["effort", "low, medium, high, max"],
            ["maxTurns", "Limit agentic turns"],
            ["memory", "user, project, local — persistent learning"],
            ["skills", "Preload skills into context"],
            ["permissionMode", "default, plan, acceptEdits, dontAsk"],
            ["isolation", "worktree — isolated git worktree"],
            ["background", "true = always run in background"],
        ])

    # --- SLIDE 26: MCP ---
    make_content_slide(prs, "Exercise 7 — MCP Servers",
        [
            "Connect Claude to external tools and services",
            "",
            ".mcp.json         — Project (shared via git)",
            "~/.claude.json    — User / local (personal)",
        ],
        code=[
            "// .mcp.json",
            '{',
            '  "mcpServers": {',
            '    "github": {',
            '      "type": "stdio",',
            '      "command": "npx",',
            '      "args": ["-y", "@anthropic-ai/github-mcp"]',
            '    }',
            '  }',
            '}',
        ])

    # --- SLIDE 27: Real-world workflow ---
    make_content_slide(prs, "Exercise 8 — Real-World Workflow",
        [
            "Putting it all together:",
            "",
            "1. Read existing code (Grep, Read)",
            "2. Plan the approach (may use Plan agent)",
            "3. Edit database.py (hook: auto-formats)",
            "4. Edit tasks.py (hook: auto-formats)",
            "5. Write tests",
            "6. Run tests (Bash)",
            "7. Commit changes",
            "",
            "CLAUDE.md conventions + hooks + multi-file coordination",
        ],
        code=[
            '"Add a search endpoint for tasks — case-insensitive',
            ' substring match on title and description via query',
            ' param \'q\'. Update database, endpoint, and tests."',
        ])

    # --- SLIDE 28: Wrap-up ---
    make_title_slide(prs, "Wrap-Up", is_section=True)

    # --- SLIDE 29: Quick wins ---
    make_content_slide(prs, "Quick Wins to Take Home",
        [
            "Today:",
            "  1. Add CLAUDE.md to your projects — biggest ROI, 5 minutes",
            "",
            "This week:",
            "  2. Configure settings — allow pytest/git, deny rm -rf/.env",
            "  3. Create one skill for your most repeated workflow",
            "",
            "This month:",
            "  4. Add hooks — auto-format, secret detection",
            "  5. Create custom agents for your team's review/research needs",
        ])

    # --- SLIDE 30: Config landscape ---
    make_table_slide(prs, "The Configuration Landscape",
        ["Feature", "File", "Guides vs Enforces"],
        [
            ["CLAUDE.md", "./CLAUDE.md", "Guides"],
            ["Settings", ".claude/settings.json", "Enforces"],
            ["Hooks", "In settings.json", "Enforces"],
            ["Skills", ".claude/skills/*/SKILL.md", "Guides"],
            ["Agents", ".claude/agents/*.md", "Both"],
            ["MCP", ".mcp.json", "Extends"],
        ])

    # --- SLIDE 31: Resources ---
    make_content_slide(prs, "Resources",
        [
            "Docs: https://code.claude.com/docs/",
            "Cheat sheet: CHEATSHEET.md in this repo",
            "Report issues: https://github.com/anthropics/claude-code/issues",
            "Exercises: exercises/ directory for self-paced review",
            "",
            "Platforms:",
            "  Terminal  |  VS Code extension  |  JetBrains plugin  |  Web app",
        ])

    # --- SLIDE 32: Questions ---
    make_title_slide(prs, "Questions?", 'claude\n> "What should I ask about Claude Code?"')

    prs.save("slides.pptx")
    print(f"Generated slides.pptx with {len(prs.slides)} slides")


if __name__ == "__main__":
    main()
